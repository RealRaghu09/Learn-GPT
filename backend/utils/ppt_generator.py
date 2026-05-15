from ast import List
from typing import Literal
from langchain_google_genai import ChatGoogleGenerativeAI
from dotenv import load_dotenv
import os
import base64
import webcolors
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import datetime
import shutil
import tempfile
import hashlib
from utils.Model import MODEL_NAME

load_dotenv()
GOOGLE_API_KEY = os.getenv('GOOGLE_API_KEY')


def _collect_downloaded_image_paths(root_dir: str) -> list[str]:
    """Collect all image file paths from the downloaded images directory."""
    image_extensions = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp'}
    image_paths = []
    
    try:
        for root, dirs, files in os.walk(root_dir):
            for file in files:
                if os.path.splitext(file)[1].lower() in image_extensions:
                    image_paths.append(os.path.join(root, file))
    except Exception as e:
        print(f"Warning: error collecting image paths: {e}")
    
    return sorted(image_paths)


def add_google_image_slides_to_presentation(
    prs: Presentation,
    keyword: str,
    max_num: int = 3,
) -> int:
    """Download images with GoogleImageCrawler, one picture per new slide, then delete the temp folder."""
    from icrawler.builtin import GoogleImageCrawler

    root_dir = tempfile.mkdtemp(prefix="learnly_ppt_images_")
    try:
        crawler = GoogleImageCrawler(storage={"root_dir": root_dir})
        crawler.crawl(keyword=keyword, max_num=max_num)
        image_paths = _collect_downloaded_image_paths(root_dir)[:max_num]
        blank_slide_layout = prs.slide_layouts[6]
        for img_path in image_paths:
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(
                img_path,
                Inches(0.75),
                Inches(1.25),
                width=Inches(8.5),
            )
        return len(image_paths)
    except Exception as e:
        print(f"Warning: image slide generation failed: {e}")
        return 0
    finally:
        shutil.rmtree(root_dir, ignore_errors=True)


TEMPERATURE  = 0.6
MAX_TOKENS = 250
class PPTModel:

    def __init__(self):
        self.warnings = []

    def generate_title_of_slides(self , topic : str , no_of_slides : int)-> list:
        try:
            llm_generate_title = ChatGoogleGenerativeAI(
                model = MODEL_NAME,
                temperature=TEMPERATURE,
                max_tokens = MAX_TOKENS,
                google_api_key=GOOGLE_API_KEY,
            )
            PROMPT_FOR_GENERATING_SLIDES_TITLE = f'''
            only Generate {no_of_slides} slided titles for the topic {topic}.
                    dont generate any other text or explanation.
            '''
            PROMPT = PROMPT_FOR_GENERATING_SLIDES_TITLE
            response =  llm_generate_title.invoke([PROMPT])
            raw_response = response.content
            list_of_topics = [line.strip() for line in raw_response.replace("*" ,"").replace("#" , "").split("\n") if line.strip()]
            if len(list_of_topics) >= no_of_slides:
                return list_of_topics[:no_of_slides]
            return list_of_topics
        except Exception as e:
            warning = f"Title generation failed: {e}"
            print(f"Warning: {warning}")
            self.warnings.append(warning)
            return [f"Slide {i + 1}" for i in range(no_of_slides)]

    def generate_content_of_topics(self ,subtopics : list[str] ,
            tone: str , depth : str ,
            no_of_slides:int ,
            TargetAudience : Literal[str],
            SlideCount : Literal[int],
            presentationStyle : Literal[str],
            PPTLang : Literal[str],
            PPTTone : Literal[str],
            SummaryLevel : Literal[str],
            PPTTheme : Literal[str]  
            ) ->str:
        llm_content = ChatGoogleGenerativeAI(
            model =MODEL_NAME,
            temperature = TEMPERATURE,
            max_tokens = 250,
            google_api_key=GOOGLE_API_KEY,
        )
        PROMPT_FOR_GENERATING_CONTENT_FOR_EACH_SLIDE = f"""
            You are an expert presentation creator.

            Generate exactly {no_of_slides} slide titles for the topic: {str(subtopics)}

            Presentation Configuration:
            - Tone: {tone}
            - Depth: {depth}
            - Number of Slides: {no_of_slides}
            - Target Audience: {TargetAudience}
            - Slide Count Preference: {SlideCount}
            - Presentation Style: {presentationStyle}
            - Presentation Language: {PPTLang}
            - Presentation Tone: {PPTTone}
            - Summary Level: {SummaryLevel}
            - Presentation Theme: {PPTTheme}

            Instructions:
            1. Generate ONLY slide titles.
            2. Do NOT generate explanations, notes, bullet points, or extra text.
            3. Titles should match the audience, tone, style, and depth provided.
            4. Ensure titles are clear, engaging, and logically ordered.
            5. Output exactly {no_of_slides} titles.
            6. Each title should be on a new line.
        """
        list_of_contents = []
        for subtopic in subtopics:
            try:
                PROMPT = PROMPT_FOR_GENERATING_CONTENT_FOR_EACH_SLIDE
                response = llm_content.invoke([PROMPT])
                raw_response = response.content
                cleaned_content = raw_response.replace("*" , "").replace("#" , "").strip()
                list_of_contents.append(cleaned_content or f"Content for {subtopic}")
            except Exception as e:
                warning = f"Content generation failed for '{subtopic}': {e}"
                print(f"Warning: {warning}")
                self.warnings.append(warning)
                list_of_contents.append(f"Content for {subtopic}")

        if not list_of_contents:
            return [f"Content for slide {i + 1}" for i in range(no_of_slides)]
        return list_of_contents


class ppt:
    def __init__(self):
        self.collected_data = []
    def extract_rgb_values(self, color: str):
        try:
            rgb_value = webcolors.name_to_rgb(color.lower())
            return list(rgb_value)
        except Exception:
            return [0, 0, 0]  # Default to black

    def add_thankyou_slide(self, prs, font, font_color, bg_color):
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # Set background color
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*bg_color)

        # Add centered "Thank You"
        box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
        tf = box.text_frame
        tf.auto_size = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "Thank You"
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.name = font
        run.font.color.rgb = RGBColor(*font_color)

    def generate_normal_ppt(self, topic: str, list_of_content: list[str], list_of_topics: list[str], no_of_slides: int,
                     color_scheme: list[str], fontstyle: list[str], include_images: bool = False) -> str:
        # Collect input data for this generation
        self.collected_data.append({
            "style": "normal",
            "topic": topic,
            "list_of_topics": list_of_topics,
            "list_of_content": list_of_content,
            "no_of_slides": no_of_slides,
            "color_scheme": color_scheme,
            "fontstyle": fontstyle,
        })
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]

        rgb1, rgb2 = [0, 0, 0], [255, 255, 255]
        if len(color_scheme) >= 2:
            rgb1 = self.extract_rgb_values(color_scheme[0])
            rgb2 = self.extract_rgb_values(color_scheme[1])

        title_font = fontstyle[0] if len(fontstyle) > 0 else "Arial"
        content_font = fontstyle[1] if len(fontstyle) > 1 else "Calibri"

        # First slide
        first_slide = prs.slides.add_slide(blank_slide_layout)
        title_shape = first_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
        tf = title_shape.text_frame
        tf.auto_size = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = topic
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.name = title_font
        run.font.color.rgb = RGBColor(*rgb1)
        first_slide.background.fill.solid()
        first_slide.background.fill.fore_color.rgb = RGBColor(*rgb2)

        # Slides 2 to no_of_slides
        for i in range(min(no_of_slides - 1, len(list_of_topics), len(list_of_content))):
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(*rgb2)

            # Title
            title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(1))
            title_tf = title_box.text_frame
            title_tf.auto_size = True
            p_title = title_tf.paragraphs[0]
            p_title.alignment = PP_ALIGN.CENTER
            run_title = p_title.add_run()
            run_title.text = list_of_topics[i]
            run_title.font.size = Pt(16)
            run_title.font.bold = True
            run_title.font.name = title_font
            run_title.font.color.rgb = RGBColor(*rgb1)

            # Content
            # New line (adds more vertical gap between title and content)
            content_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(5))
            content_tf = content_box.text_frame
            content_tf.auto_size = True
            p_content = content_tf.paragraphs[0]
            p_content.alignment = PP_ALIGN.CENTER
            run_content = p_content.add_run()
            run_content.text = list_of_content[i]
            run_content.font.size = Pt(14)
            run_content.font.name = content_font
            run_content.font.color.rgb = RGBColor(*rgb1)

        if include_images:
            try:
                image_kw = (topic[:120] if topic else "presentation").strip() or "presentation"
                add_google_image_slides_to_presentation(
                    prs,
                    keyword=image_kw,
                    max_num=min(5, max(1, no_of_slides - 1)),
                )
            except Exception as e:
                print(f"Warning: skipped images for normal PPT because: {e}")

        # ➕ Adding "Thank You" slide at end
        self.add_thankyou_slide(prs, font=title_font, font_color=tuple(rgb1), bg_color=tuple(rgb2))

        # Save with hashed filename
        timestamp = datetime.datetime.now().isoformat()
        unique_hash = hashlib.md5((topic + timestamp).encode()).hexdigest()
        output_dir = os.path.abspath("generated_ppt")
        os.makedirs(output_dir, exist_ok=True)
        filename = os.path.join(output_dir, f"generated_presentation_{unique_hash}.pptx")
        prs.save(filename)

        # Clear transient data after generation completes
        self.collected_data.clear()

        return {"message": f"Generated Normal Successfully! Saved as {filename}", "file_path": filename}
    

    def generate_modern_ppt(self, topic: str, list_of_content: list[str], list_of_topics: list[str], no_of_slides: int,
                            include_images: bool = False) -> str:
        # Collect input data for this generation
        self.collected_data.append({
            "style": "modern",
            "topic": topic,
            "list_of_topics": list_of_topics,
            "list_of_content": list_of_content,
            "no_of_slides": no_of_slides,
        })

        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]

        # Color palette: title color (accent), bg color (light/dark)
        rgb1 = [30, 30, 30]   # Dark Gray/Black
        rgb2 = [245, 245, 245] # Light Gray

        title_font = "Montserrat"
        content_font =  "Lato"

        # First slide: Title
        first_slide = prs.slides.add_slide(blank_slide_layout)
        first_slide.background.fill.solid()
        first_slide.background.fill.fore_color.rgb = RGBColor(*rgb2)

        title_box = first_slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
        tf = title_box.text_frame
        tf.auto_size = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = topic
        run.font.size = Pt(44)
        run.font.bold = True
        run.font.name = title_font
        run.font.color.rgb = RGBColor(*rgb1)

        # Content slides
        for i in range(min(no_of_slides - 1, len(list_of_topics), len(list_of_content))):
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(*rgb2)

            # Title
            title_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(1))
            title_tf = title_box.text_frame
            title_tf.auto_size = True
            p_title = title_tf.paragraphs[0]
            p_title.alignment = PP_ALIGN.LEFT
            run_title = p_title.add_run()
            run_title.text = list_of_topics[i]
            run_title.font.size = Pt(28)
            run_title.font.bold = True
            run_title.font.name = title_font
            run_title.font.color.rgb = RGBColor(*rgb1)

            # Content
            content_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8.5), Inches(5))
            content_tf = content_box.text_frame
            content_tf.word_wrap = True
            content_tf.auto_size = True
            p_content = content_tf.paragraphs[0]
            p_content.alignment = PP_ALIGN.LEFT
            run_content = p_content.add_run()
            run_content.text = list_of_content[i]
            run_content.font.size = Pt(18)
            run_content.font.name = content_font
            run_content.font.color.rgb = RGBColor(*rgb1)

        if include_images:
            try:
                image_kw = (topic[:120] if topic else "presentation").strip() or "presentation"
                add_google_image_slides_to_presentation(
                    prs,
                    keyword=image_kw,
                    max_num=min(5, max(1, no_of_slides - 1)),
                )
            except Exception as e:
                print(f"Warning: skipped images for modern PPT because: {e}")

        # Add modern-style thank you slide
        self.add_thankyou_slide(prs, font=title_font, font_color=tuple(rgb1), bg_color=tuple(rgb2))

        # Save presentation
        timestamp = datetime.datetime.now().isoformat()
        unique_hash = hashlib.md5((topic + timestamp).encode()).hexdigest()
        output_dir = os.path.abspath("generated_ppt")
        os.makedirs(output_dir, exist_ok=True)
        filename = os.path.join(output_dir, f"modern_presentation_{unique_hash}.pptx")
        prs.save(filename)
        
        # Clear transient data after generation completes
        self.collected_data.clear()

        return {"message":f"Modern PPT Generated Successfully! Saved as {filename}" , "file_path" : filename}

    def generate_creative_ppt(self, topic: str, list_of_content: list[str], list_of_topics: list[str], no_of_slides: int,
                              include_images: bool = False) -> str:
        # Collect input data for this generation
        self.collected_data.append({
            "style": "creative",
            "topic": topic,
            "list_of_topics": list_of_topics,
            "list_of_content": list_of_content,
            "no_of_slides": no_of_slides,
        })
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        import datetime, os

        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]

        # Colors
        rgb1 =  [255, 255, 255]  # Text
        rgb2 =  [0, 0, 0]        # Background

        title_font =  "Poppins"
        content_font = "Comic Sans MS"

        # First slide (Title slide - center aligned and bold)
        first_slide = prs.slides.add_slide(blank_slide_layout)
        first_slide.background.fill.solid()
        first_slide.background.fill.fore_color.rgb = RGBColor(*rgb2)

        title_box = first_slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2.5))
        tf = title_box.text_frame
        tf.auto_size = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = topic.upper()
        run.font.size = Pt(48)
        run.font.bold = True
        run.font.name = title_font
        run.font.color.rgb = RGBColor(*rgb1)

        # Content slides
        for i in range(min(no_of_slides - 1, len(list_of_topics), len(list_of_content))):
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(*rgb2)

            # Title (bold and colored)
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.5), Inches(1))
            title_tf = title_box.text_frame
            title_tf.auto_size = True
            p_title = title_tf.paragraphs[0]
            p_title.alignment = PP_ALIGN.LEFT
            run_title = p_title.add_run()
            run_title.text = f"🎨 {list_of_topics[i]}"
            run_title.font.size = Pt(32)
            run_title.font.bold = True
            run_title.font.name = title_font
            run_title.font.color.rgb = RGBColor(*rgb1)

            # Content (centered and expressive)
            content_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8.5), Inches(5))
            content_tf = content_box.text_frame
            content_tf.word_wrap = True
            content_tf.auto_size = True
            p_content = content_tf.paragraphs[0]
            p_content.alignment = PP_ALIGN.LEFT
            run_content = p_content.add_run()
            run_content.text = list_of_content[i]
            run_content.font.size = Pt(20)
            run_content.font.name = content_font
            run_content.font.color.rgb = RGBColor(*rgb1)

        if include_images:
            try:
                image_kw = (topic[:120] if topic else "presentation").strip() or "presentation"
                add_google_image_slides_to_presentation(
                    prs,
                    keyword=image_kw,
                    max_num=min(5, max(1, no_of_slides - 1)),
                )
            except Exception as e:
                print(f"Warning: skipped images for creative PPT because: {e}")

        # Add creative thank-you slide
        self.add_thankyou_slide(prs, font=title_font, font_color=tuple(rgb1), bg_color=tuple(rgb2))

        # Save file
        timestamp = datetime.datetime.now().isoformat()
        unique_hash = hashlib.md5((topic + timestamp).encode()).hexdigest()
        output_dir = os.path.abspath("generated_ppt")
        os.makedirs(output_dir, exist_ok=True)
        filename = os.path.join(output_dir, f"creative_presentation_{unique_hash}.pptx")
        prs.save(filename)
        
        # Clear transient data after generation completes
        self.collected_data.clear()

        return {"message" : f"Creative PPT Generated Successfully! Saved as {filename}" , "file_path" : filename}
    def generate_retro_ppt(self, topic: str, list_of_content: list[str], list_of_topics: list[str], no_of_slides: int,
                           include_images: bool = False) -> str:
        # Collect input data for this generation
        self.collected_data.append({
            "style": "retro",
            "topic": topic,
            "list_of_topics": list_of_topics,
            "list_of_content": list_of_content,
            "no_of_slides": no_of_slides,
        })

        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]

        # Retro default palette (fallbacks)
        rgb1 = [139, 69, 19]   # Saddle Brown
        rgb2 = [255, 228, 181] # Moccasin

        title_font = "Courier New"
        content_font = "Georgia"

        # First slide - Title
        first_slide = prs.slides.add_slide(blank_slide_layout)
        first_slide.background.fill.solid()
        first_slide.background.fill.fore_color.rgb = RGBColor(*rgb2)

        title_box = first_slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2.5))
        tf = title_box.text_frame
        tf.auto_size = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = topic.upper()
        run.font.size = Pt(42)
        run.font.bold = True
        run.font.name = title_font
        run.font.color.rgb = RGBColor(*rgb1)

        # Content slides
        for i in range(min(no_of_slides - 1, len(list_of_topics), len(list_of_content))):
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(*rgb2)

            # Title
            title_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(1))
            title_tf = title_box.text_frame
            title_tf.auto_size = True
            p_title = title_tf.paragraphs[0]
            p_title.alignment = PP_ALIGN.LEFT
            run_title = p_title.add_run()
            run_title.text = f"★ {list_of_topics[i]}"
            run_title.font.size = Pt(28)
            run_title.font.bold = True
            run_title.font.name = title_font
            run_title.font.color.rgb = RGBColor(*rgb1)

            # Content
            content_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8.5), Inches(5))
            content_tf = content_box.text_frame
            content_tf.word_wrap = True
            content_tf.auto_size = True
            p_content = content_tf.paragraphs[0]
            p_content.alignment = PP_ALIGN.LEFT
            run_content = p_content.add_run()
            run_content.text = list_of_content[i]
            run_content.font.size = Pt(18)
            run_content.font.name = content_font
            run_content.font.color.rgb = RGBColor(*rgb1)

        if include_images:
            try:
                image_kw = (topic[:120] if topic else "presentation").strip() or "presentation"
                add_google_image_slides_to_presentation(
                    prs,
                    keyword=image_kw,
                    max_num=min(5, max(1, no_of_slides - 1)),
                )
            except Exception as e:
                print(f"Warning: skipped images for retro PPT because: {e}")

        self.add_thankyou_slide(prs, font=title_font, font_color=tuple(rgb1), bg_color=tuple(rgb2))

        # Save
        timestamp = datetime.datetime.now().isoformat()
        unique_hash = hashlib.md5((topic + timestamp).encode()).hexdigest()
        output_dir = os.path.abspath("generated_ppt")
        os.makedirs(output_dir, exist_ok=True)
        filename = os.path.join(output_dir, f"retro_presentation_{unique_hash}.pptx")
        prs.save(filename)
        
        # Clear transient data after generation completes
        self.collected_data.clear()

        return {"message":f"Retro PPT Generated Successfully! Saved as {filename}" , "file_path" : filename}

def get_link(file_path : str):
    with open(file_path, "rb") as f:
        ppt_bytes = f.read()

    encoded = base64.b64encode(ppt_bytes).decode('utf-8')

    # Raw data URI link (safe to use as href without tags)
    download_link = f"data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{encoded}"

    return download_link